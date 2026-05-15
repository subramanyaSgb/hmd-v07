"""
Build the HMD project status presentation as a native .pptx file.

Generates a 33-slide deck covering:
  A. Project Foundation (4 slides)
  B. Project Timeline (3)
  C. Infrastructure & Server Setup (3)
  D. Data Sources & Integration Journey (4)
  E. Application Capabilities — V7 Proper (6)
  F. Current Operations Status (3)
  G. Data Reality / Gap Findings (4)
  H. Clarifications Needed from JSW (3)
  I. Path Forward (2)
  J. Close (1)

Run:
    python _build_status_pptx.py
Produces:
    2026-05-14-hmd-project-status.pptx in the same directory.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree


# ─── Design system ────────────────────────────────────────────────

# Primary palette — corporate, restrained
NAVY      = RGBColor(0x1E, 0x3A, 0x8A)    # dark navy — primary
NAVY_LITE = RGBColor(0x3B, 0x5B, 0xA9)    # lighter navy — accents
GOLD      = RGBColor(0xC2, 0x8B, 0x1B)    # gold — highlights / status: at-risk
GREEN     = RGBColor(0x15, 0x80, 0x3D)    # status: complete / on-track
AMBER     = RGBColor(0xB4, 0x53, 0x09)    # status: delayed / partial
RED       = RGBColor(0xB9, 0x1C, 0x1C)    # status: blocker / overdue
GRAY_DARK = RGBColor(0x1F, 0x29, 0x37)    # body text
GRAY_MED  = RGBColor(0x4B, 0x55, 0x63)    # secondary text
GRAY_LITE = RGBColor(0xE5, 0xE7, 0xEB)    # light backgrounds
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
LINE      = RGBColor(0xCB, 0xD5, 0xE1)

# Font families
TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"
MONO_FONT  = "Consolas"

# Slide dimensions — 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Layout margins
MARGIN_L = Inches(0.5)
MARGIN_R = Inches(0.5)
MARGIN_T = Inches(0.4)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
CONTENT_H = SLIDE_H - MARGIN_T - Inches(0.4)


# ─── Helpers ──────────────────────────────────────────────────────

def add_solid_rect(slide, left, top, width, height, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def add_text(slide, left, top, width, height, text, *,
             font_name=BODY_FONT, font_size=14, bold=False, italic=False,
             color=GRAY_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             word_wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        text = [text]

    for i, line in enumerate(text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                font_size=14, color=GRAY_DARK, line_spacing=1.15,
                bullet_color=None):
    """items: list of strings, or list of (str, level) tuples for sub-bullets."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.name = BODY_FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        # add bullet point manually via XML — python-pptx default bullets are
        # inconsistent
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        # Remove any existing bullet props
        for existing in pPr.findall(qn('a:buChar')) + pPr.findall(qn('a:buNone')):
            pPr.remove(existing)
        # Add bullet character
        if level == 0:
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '■')
        else:
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '–')
    return tb


def add_table(slide, left, top, width, height, headers, rows, *,
              header_bg=NAVY, header_fg=WHITE, body_fg=GRAY_DARK,
              alt_row_bg=OFF_WHITE, font_size=11, header_size=12,
              col_widths=None):
    """Add a styled table. rows is a list of lists; first list = header row."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        run.font.name = BODY_FONT
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = header_fg

    # Data rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE if i % 2 == 0 else WHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.03)
            tf.margin_bottom = Inches(0.03)
            tf.word_wrap = True
            # Each value can be a (text, color, bold) tuple for highlighting
            if isinstance(val, tuple):
                text, fg, bold = val
            else:
                text, fg, bold = val, body_fg, False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(text)
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = fg
            run.font.bold = bold

    return table


def add_slide_header(slide, section_letter, section_name, slide_title,
                     slide_number=None):
    """Standard slide header: section strip + slide title."""
    # Section strip (top-left small navy bar)
    add_solid_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, NAVY)

    # Section label (rotated text in the strip)
    # python-pptx doesn't easily rotate text; instead show it across a top-right strip
    add_text(slide, Inches(0.55), Inches(0.25), Inches(7), Inches(0.32),
             f"SECTION {section_letter} · {section_name}",
             font_size=9, bold=True, color=NAVY_LITE)

    # Slide title
    add_text(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.7),
             slide_title,
             font_size=26, bold=True, color=NAVY)

    # Thin separator line
    add_solid_rect(slide, Inches(0.55), Inches(1.35),
                   Inches(12.2), Emu(11430),     # ~0.0125 inch = thin line
                   NAVY)

    # Footer
    add_text(slide, Inches(0.55), Inches(7.1), Inches(8), Inches(0.3),
             "HMD Project Status — Deevia Software for JSW Vijayanagar — 14 May 2026",
             font_size=9, color=GRAY_MED)

    if slide_number is not None:
        add_text(slide, Inches(12.4), Inches(7.1), Inches(0.6), Inches(0.3),
                 f"{slide_number}", font_size=9, color=GRAY_MED,
                 align=PP_ALIGN.RIGHT)


def add_section_divider(slide, section_letter, section_name, subtitle=""):
    """Full-bleed section divider slide."""
    add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)

    # Big "SECTION X"
    add_text(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.6),
             f"SECTION {section_letter}",
             font_size=18, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Section name
    add_text(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.5),
             section_name,
             font_size=46, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        add_text(slide, Inches(1), Inches(4.6), Inches(11.3), Inches(0.7),
                 subtitle,
                 font_size=18, italic=True, color=RGBColor(0xC7, 0xD2, 0xFE),
                 align=PP_ALIGN.LEFT)


def add_cover_slide(slide):
    """Slide 1 — cover."""
    # Background
    add_solid_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)

    # Gold accent bar
    add_solid_rect(slide, Inches(0), Inches(2.0), Inches(1.5), Inches(0.1), GOLD)

    # Project name
    add_text(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(0.6),
             "HMD PROJECT — STATUS PRESENTATION",
             font_size=14, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Title
    add_text(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6),
             "Hot Metal Distribution System",
             font_size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.7),
             "End-to-end project review — design, development, integration, current state, path forward",
             font_size=18, italic=True, color=RGBColor(0xC7, 0xD2, 0xFE),
             align=PP_ALIGN.LEFT)

    # Footer block
    add_solid_rect(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(1.0), NAVY_LITE)
    add_text(slide, Inches(1), Inches(6.65), Inches(6), Inches(0.4),
             "PREPARED FOR",
             font_size=9, bold=True, color=GOLD)
    add_text(slide, Inches(1), Inches(6.95), Inches(6), Inches(0.4),
             "JSW Steel Limited — Vijayanagar Works",
             font_size=15, bold=True, color=WHITE)
    add_text(slide, Inches(8), Inches(6.65), Inches(5), Inches(0.4),
             "PREPARED BY",
             font_size=9, bold=True, color=GOLD)
    add_text(slide, Inches(8), Inches(6.95), Inches(5), Inches(0.4),
             "Deevia Software (India) Pvt. Ltd. — 14 May 2026",
             font_size=15, bold=True, color=WHITE)


# ─── Slide builders ───────────────────────────────────────────────

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])     # blank
    add_cover_slide(s)


def slide_agenda(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "—", "OVERVIEW", "Agenda", slide_number=n)

    items_left = [
        ("A.  Project Foundation", 0),
        ("Goal, scope, system architecture", 1),
        ("", 0),
        ("B.  Project Timeline", 0),
        ("Original plan vs actual progress", 1),
        ("", 0),
        ("C.  Infrastructure & Server Setup", 0),
        ("Hardware, network, software stack", 1),
        ("", 0),
        ("D.  Data Sources & Integration", 0),
        ("Three external sources, integration dates", 1),
    ]
    items_right = [
        ("E.  Application Capabilities — V7 Proper", 0),
        ("Module-by-module walkthrough", 1),
        ("", 0),
        ("F.  Current Operations Status", 0),
        ("What's live, what's verified", 1),
        ("", 0),
        ("G.  Data Reality / Findings", 0),
        ("What's reliable, what isn't, key gaps", 1),
        ("", 0),
        ("H.  Clarifications Needed from JSW", 0),
        ("", 0),
        ("I.  Path Forward + J.  Close", 0),
    ]
    add_bullets(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(5.4),
                items_left, font_size=14)
    add_bullets(s, Inches(6.9), Inches(1.6), Inches(6.0), Inches(5.4),
                items_right, font_size=14)


def slide_project_goal(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "A", "PROJECT FOUNDATION", "Project goal and scope", slide_number=n)

    # Left column — the problem
    add_text(s, Inches(0.7), Inches(1.6), Inches(5.8), Inches(0.4),
             "THE OPERATIONAL CHALLENGE",
             font_size=11, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.8),
             ["JSW Vijayanagar Works operates a hot metal distribution network connecting 7 producers (BF1–BF5, COREX1, COREX2) to 4 steel melting shops (SMS-1 through SMS-4) using a fleet of 53 torpedo ladles.",
              "",
              "Coordinating torpedo allocation, trip planning, real-time tracking, and per-trip lifecycle management across the network is operationally complex and currently relies on disparate sources of data without a unified planning or monitoring layer."],
             font_size=14)

    # Vertical divider
    add_solid_rect(s, Inches(6.65), Inches(1.6), Emu(11430), Inches(5.0), LINE)

    # Right column — the solution
    add_text(s, Inches(6.95), Inches(1.6), Inches(5.8), Inches(0.4),
             "WHAT HMD DELIVERS",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(6.95), Inches(2.0), Inches(5.8), Inches(5.0), [
        "Centralised real-time view of every torpedo and every trip",
        "Role-based access for Admin, Producer, Consumer operators",
        "Daily and monthly planning with optimisation (PuLP) for route assignment",
        "Trip lifecycle management — 16 distinct states from creation to completion",
        "Per-node operations monitoring, deviation analytics, and KPI dashboards",
        "Audit trail, alerts, reporting (PDF/Excel), and maintenance scheduling",
        "Integration with JSW source systems — torpedo GPS, BF weighbridge transactions, SMS hot-metal receipts",
    ], font_size=12.5)


def slide_architecture(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "A", "PROJECT FOUNDATION", "System architecture", slide_number=n)

    # Title strip
    add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
             "HMD ARCHITECTURE — A SINGLE WEB APPLICATION INTEGRATING THREE JSW DATA SOURCES",
             font_size=10, bold=True, color=GOLD)

    # User layer (top)
    add_solid_rect(s, Inches(1.0), Inches(2.05), Inches(11.3), Inches(0.7), NAVY_LITE)
    add_text(s, Inches(1.0), Inches(2.05), Inches(11.3), Inches(0.7),
             "USERS — JSW operators (Admin / Producer / Consumer roles) via web browser over JSW network",
             font_size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Arrow down (simulated with a thin rectangle)
    add_solid_rect(s, Inches(6.55), Inches(2.85), Inches(0.2), Inches(0.3), GRAY_MED)

    # HMD application layer
    add_solid_rect(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(1.3),
                   RGBColor(0xEF, 0xF6, 0xFF))
    add_solid_rect(s, Inches(1.0), Inches(3.2), Inches(0.15), Inches(1.3), NAVY)
    add_text(s, Inches(1.25), Inches(3.25), Inches(11), Inches(0.35),
             "HMD APPLICATION — running on BF4 PC (JSW network)",
             font_size=11, bold=True, color=NAVY)
    add_text(s, Inches(1.25), Inches(3.6), Inches(11), Inches(0.85),
             ["• Frontend (React + Vite) — Dashboard, Trip Mgmt, Planning, Operations, Statistics, Reports, Audit, Fleet, Maintenance, Settings",
              "• Backend (FastAPI + Python) — REST API, RBAC, JWT auth, audit trail, alert detector, optimisation engine, scheduled syncs (APScheduler)",
              "• Local storage (PostgreSQL + Redis) — application state, mirrors of upstream data, fleet & trip records, configs"],
             font_size=10.5)

    # Arrow down
    add_solid_rect(s, Inches(6.55), Inches(4.6), Inches(0.2), Inches(0.3), GRAY_MED)

    # Data sources layer
    db_box_w = Inches(3.5)
    db_box_h = Inches(1.7)
    db_y = Inches(4.95)
    db_xs = [Inches(1.0), Inches(4.92), Inches(8.83)]

    sources = [
        ("SuVeechi MySQL", "Real-time torpedo GPS",
         "• vw_unit_status_ist\n• 53 torpedoes\n• 10-second refresh\n• Live since 06-May-2026"),
        ("WBATNGL Oracle", "BF weighbridge transactions",
         "• BF3 + BF5 schemas\n• 8 of 9 views working\n• 60-second sync\n• Trip mirror since 08-May-2026"),
        ("HTS Oracle", "SMS hot-metal heats",
         "• VW_HTS_HOTMETAL_DATA\n• JVMLPROD instance\n• 5-minute sync\n• Access set up 11-May-2026"),
    ]
    for x, (title, sub, body) in zip(db_xs, sources):
        add_solid_rect(s, x, db_y, db_box_w, db_box_h, OFF_WHITE)
        add_solid_rect(s, x, db_y, db_box_w, Inches(0.45), NAVY)
        add_text(s, x + Inches(0.1), db_y, db_box_w - Inches(0.2), Inches(0.45),
                 title,
                 font_size=12, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), db_y + Inches(0.5), db_box_w - Inches(0.2), Inches(0.3),
                 sub,
                 font_size=10, italic=True, color=NAVY_LITE)
        add_text(s, x + Inches(0.1), db_y + Inches(0.8), db_box_w - Inches(0.2), db_box_h - Inches(0.85),
                 body, font_size=10, color=GRAY_DARK)


# ─── SECTION B — Timeline ─────────────────────────────────────────

def slide_section_b(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(s, "B", "Project Timeline",
                        "Original plan, actual progress, and where we stand today.")


def slide_b_timeline_plan(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "B", "PROJECT TIMELINE",
                     "Original plan — five milestones, Dec 2025 to Aug 2026", slide_number=n)

    rows = [
        ["M1", "Requirements Collection & Design", "08-Dec-2025", "15-Jan-2026", ("Complete", GREEN, True)],
        ["M2", "Core HMD Development", "15-Jan-2026", "22-May-2026", ("98% — Active", GREEN, True)],
        ["M3", "Dashboard & Integration", "15-Apr-2026", "19-Jun-2026", ("65% — Active", AMBER, True)],
        ["M4", "Testing Simulation & Pilot", "26-Jun-2026", "24-Jul-2026", ("Not Started", GRAY_MED, False)],
        ["M5", "Full Go-Live & Stabilization", "24-Jul-2026", "07-Aug-2026", ("Not Started", GRAY_MED, False)],
    ]
    add_table(s, Inches(0.7), Inches(1.7), Inches(12), Inches(2.7),
              ["ID", "Milestone", "Start Date", "End Date", "Status as of 14-May-2026"],
              rows,
              col_widths=[Inches(0.7), Inches(4.5), Inches(2.2), Inches(2.2), Inches(2.4)],
              font_size=13, header_size=13)

    # Note below
    add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
             "KEY POINTS",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2.0), [
        "M1 closed on schedule — requirements, hardware spec, queries baseline established.",
        "M2 is 98% complete — the core HMD application has shipped; remaining 2% reserved for refinements after live-data observation.",
        "M3 (Dashboard & Integration) was blocked until April 2026 due to network and database access; now actively progressing after VPN granted 06-May-2026.",
        "M4 and M5 are downstream of M3 — start dates were shifted +1 week as per Weekly Report 12 (23-Apr-2026).",
    ], font_size=13)


def slide_b_timeline_actual(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "B", "PROJECT TIMELINE",
                     "What actually happened — major events on the journey",
                     slide_number=n)

    events = [
        ["Dec-2025", "Project kickoff — Queries shared with JSW, hardware spec submitted (HPE ProLiant DL160 Gen9 + Windows Server 2022)"],
        ["Jan-2026", "Historical data received from JSW — 1-year torpedo loading, cycle time report, weighbridge report, maintenance reports"],
        ["Jan–Apr 2026", "Core development — 216 backend + frontend tasks completed across all V7 modules"],
        ["Feb-2026", "Windows Server 2022 license procured"],
        ["Mar-Apr 2026", "Data requirements iteration with JSW IT — 3 source-system access requests, multiple Teams meetings"],
        ["02-Apr-2026", "SuVeechi MySQL view shared by vendor (Ganesha)"],
        ["08-Apr-2026", "Network access escalation — 100% packet loss to JSW DB servers; root cause traced to absent LAN"],
        ["14-Apr-2026", "Formal approvals closed — Gagan Chopra approval forwarded; DEP008 closed"],
        ["~30-Apr-2026", "HMD PC physically integrated into JSW network — placed at BF4"],
        ["06-May-2026", "VPN access granted to Deevia team — DEP009 + DEP012 closed; SuVeechi sync went LIVE — 53 torpedoes streaming GPS every 10s"],
        ["07-May-2026", "Live Tracking Phase 1–3 shipped — map polish, Torpedo Drawer, per-torpedo capacity backfill"],
        ["08-May-2026", "WBATNGL Trip Mirror sprint shipped — 60-second sync, BF-side trip data flowing into HMD"],
        ["11-May-2026", "Oracle HTS connectivity unblocked — service name + user identified; DEP007 closed"],
    ]
    rows = []
    for d, desc in events:
        rows.append([(d, NAVY, True), desc])
    add_table(s, Inches(0.7), Inches(1.6), Inches(12), Inches(5.3),
              ["Date", "Event"],
              rows,
              col_widths=[Inches(1.9), Inches(10.1)],
              font_size=11, header_size=12)


def slide_b_milestone_progress(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "B", "PROJECT TIMELINE",
                     "Current milestone position — M2 nearly closed, M3 active, M4-M5 ahead",
                     slide_number=n)

    rows = [
        ["M1", "Requirements & Design",         "100%",         "100%",        ("Complete", GREEN, True)],
        ["M2", "Core Development",              "100% (planned)","98%",         ("On Track", GREEN, True)],
        ["M3", "Dashboard & Integration",       "65% (planned)","65%",         ("On Track — major jump in May", AMBER, True)],
        ["M4", "Testing & Pilot",               "0%",           "0%",          ("Awaiting M3 completion", GRAY_MED, False)],
        ["M5", "Go-Live & Stabilization",       "0%",           "0%",          ("Awaiting M4 completion", GRAY_MED, False)],
    ]
    add_table(s, Inches(0.7), Inches(1.6), Inches(12), Inches(2.5),
              ["ID", "Milestone", "Planned", "Actual", "Comment"],
              rows,
              col_widths=[Inches(0.7), Inches(4.2), Inches(1.6), Inches(1.6), Inches(3.9)],
              font_size=13, header_size=13)

    add_text(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.4),
             "OVERALL POSITION (14-May-2026)",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.7), Inches(4.6), Inches(12), Inches(2.4), [
        "Overall project progress: approximately 53% complete against the planned timeline.",
        "M3 made a substantial jump in early May (from ~18% to 65%) as VPN access enabled remote testing of all three source databases.",
        "Remaining 35% of M3 depends on: (a) Level 1/2 production signals (DEP002, overdue), (b) Production/Consumption rules documentation (DEP005, overdue), (c) end-to-end pilot validation against real plant operations.",
        "M4 (Testing) and M5 (Go-Live) are tracking start dates of 26-Jun-2026 and 24-Jul-2026 respectively, contingent on M3 completion.",
    ], font_size=13)


# ─── SECTION C — Infrastructure ───────────────────────────────────

def slide_section_c(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(s, "C", "Infrastructure & Server Setup",
                        "Hardware, network, and software stack for HMD at JSW Vijayanagar.")


def slide_c_server(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "C", "INFRASTRUCTURE & SERVER SETUP",
                     "Server topology and physical deployment",
                     slide_number=n)

    # Left half — server specs
    add_text(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.4),
             "HMD HOST SERVER (AT JSW BF4 LOCATION)",
             font_size=11, bold=True, color=GOLD)
    rows = [
        [("Hardware", NAVY, True),  "HPE ProLiant DL160 Gen9"],
        [("Operating System", NAVY, True), "Windows Server 2022 (license procured Feb-2026)"],
        [("Application Stack", NAVY, True), "FastAPI · React · PostgreSQL · Redis"],
        [("Runtime", NAVY, True), "Python venv + Node.js + nginx (frontend)"],
        [("Physical Location", NAVY, True), "BF4 (Blast Furnace 4 area)"],
        [("Owned & Maintained by", NAVY, True), "JSW BF4 team (per 22-Apr-2026 agreement)"],
    ]
    add_table(s, Inches(0.7), Inches(1.85), Inches(6), Inches(3.2),
              ["Specification", "Detail"],
              rows,
              col_widths=[Inches(2.0), Inches(4.0)],
              font_size=11, header_size=11)

    # Right half — access details
    add_text(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4),
             "ACCESS & CONNECTIVITY",
             font_size=11, bold=True, color=GOLD)
    rows = [
        [("Frontend access", NAVY, True), "http://<BF4-PC>:5173 (JSW intranet)"],
        [("Backend API", NAVY, True), "http://<BF4-PC>:8000"],
        [("Database (internal)", NAVY, True), "PostgreSQL on localhost:5432"],
        [("Cache (internal)", NAVY, True), "Redis on localhost:6379"],
        [("Remote development", NAVY, True), "VPN to JSW network → AnyDesk to BF4 PC"],
        [("Source DB access", NAVY, True), "Outbound to JSW DB servers (SuVeechi 3306, Oracle 1522)"],
    ]
    add_table(s, Inches(7.0), Inches(1.85), Inches(6), Inches(3.2),
              ["Aspect", "Detail"],
              rows,
              col_widths=[Inches(2.0), Inches(4.0)],
              font_size=11, header_size=11)

    add_text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(0.4),
             "KEY POINT",
             font_size=11, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.4),
             "The HMD application lives entirely inside JSW's secure network. No data leaves the plant. All three source databases (SuVeechi, WBATNGL, HTS) are reached from the BF4 PC over JSW's internal network. Deevia developers connect via approved VPN credentials (granted 06-May-2026) for ongoing support and updates only — no production data flows outside JSW infrastructure.",
             font_size=13)


def slide_c_network(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "C", "INFRASTRUCTURE & SERVER SETUP",
                     "Network integration journey",
                     slide_number=n)

    events = [
        ["Mar-2026", "Initial network plan — Deevia dev system to be placed at BF4-L2 subnet (192.168.150.100)"],
        ["08-Apr-2026", "Network access request raised — DB servers unreachable from Deevia PC (100% packet loss)"],
        ["17-Apr-2026", "Root cause confirmed — no LAN connection from DSI laptop (firewall ruled out)"],
        ["21-Apr-2026", "Teams meeting with JSW IT (Divakara, Hari, Pradeep, Vasagerappa) to resolve LAN integration approach"],
        ["22-Apr-2026", "Server hardware maintenance ownership agreed with BF4 team (Divakar's email)"],
        ["~24–30-Apr-2026", "PC physically delivered to JSW IT, integrated into BF4 network, returned operational"],
        ["~30-Apr-2026", "BF4 LAN connectivity verified — all 3 DB servers reachable from BF4 PC"],
        ["06-May-2026", "VPN access for Deevia team granted (Vamsi Krishna obtained domain ID from 365 Support; Raghavendra/JSW IT configured VPN). DEP009 + DEP012 both closed."],
    ]
    rows = []
    for d, desc in events:
        rows.append([(d, NAVY, True), desc])
    add_table(s, Inches(0.7), Inches(1.6), Inches(12), Inches(4.0),
              ["Date", "Network Integration Event"],
              rows,
              col_widths=[Inches(2.0), Inches(10.0)],
              font_size=12, header_size=13)

    add_text(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.4),
             "OUTCOME",
             font_size=11, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.9),
             "By 06-May-2026, the BF4 PC was reachable from Deevia development laptops over JSW VPN, and outbound connectivity to all three source databases was operational. This was the single biggest unblocking event of M3 — within two days of access, both SuVeechi and WBATNGL syncs went live.",
             font_size=13)


def slide_c_stack(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "C", "INFRASTRUCTURE & SERVER SETUP",
                     "Software stack and security posture",
                     slide_number=n)

    # Stack tiers
    add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
             "APPLICATION STACK — STANDARD ENTERPRISE WEB APPLICATION",
             font_size=11, bold=True, color=GOLD)
    tiers = [
        ("PRESENTATION", NAVY, "React 19 + Vite 7 + React Router v7 — single-page web app, role-based UI, TypeScript-incremental.\nCharts: Recharts.  Maps: Leaflet.  PDF/Excel export: jsPDF, xlsx."),
        ("APPLICATION", NAVY_LITE, "FastAPI + Uvicorn — async REST API. 13 SQLAlchemy ORM models. APScheduler for background syncs.\nOptimisation: PuLP (linear programming for daily plan generation).  Tracing: OpenTelemetry."),
        ("DATA",        GOLD,     "PostgreSQL (primary). Redis (cache with in-memory fallback). Alembic-managed schema migrations.\nLocal mirrors of SuVeechi / WBATNGL / HTS for resilient analytics."),
        ("SECURITY",    AMBER,    "JWT authentication. bcrypt password hashing. RBAC (Admin / Producer / Consumer).\nAccount lockout, rate limiting, CSRF protection, audit trail with username + IP."),
    ]
    y = Inches(1.85)
    for label, color, body in tiers:
        # Tier label box
        add_solid_rect(s, Inches(0.7), y, Inches(1.8), Inches(1.05), color)
        add_text(s, Inches(0.7), y, Inches(1.8), Inches(1.05),
                 label,
                 font_size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body box
        add_solid_rect(s, Inches(2.55), y, Inches(10.2), Inches(1.05), OFF_WHITE)
        add_text(s, Inches(2.65), y + Inches(0.1), Inches(10.0), Inches(0.95),
                 body,
                 font_size=11, color=GRAY_DARK)
        y += Inches(1.2)


# ─── SECTION D — Data Sources ─────────────────────────────────────

def slide_section_d(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(s, "D", "Data Sources & Integration Journey",
                        "Three external sources of truth, integrated into HMD over April–May 2026.")


def slide_d_sources_overview(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "D", "DATA SOURCES & INTEGRATION",
                     "Three external sources — at a glance",
                     slide_number=n)

    rows = [
        [("SuVeechi", NAVY, True), "MySQL", "vw_unit_status_ist", "53 torpedoes — live GPS, status, location", "10 seconds", ("Live since 06-May", GREEN, True)],
        [("WBATNGL",  NAVY, True), "Oracle", "BF3 + BF5 schemas (9 views)", "BF-side weighbridge transactions, tap data, chemistry", "60 seconds", ("Live since 08-May", GREEN, True)],
        [("HTS",      NAVY, True), "Oracle", "VW_HTS_HOTMETAL_DATA + 4 tables", "SMS hot-metal receipts, caster events (SMS-4)", "5 minutes", ("Access set up 11-May", AMBER, True)],
    ]
    add_table(s, Inches(0.55), Inches(1.6), Inches(12.4), Inches(2.4),
              ["Source", "Engine", "Primary view", "What it provides", "Sync rate", "Status"],
              rows,
              col_widths=[Inches(1.2), Inches(1.0), Inches(2.6), Inches(4.5), Inches(1.4), Inches(1.7)],
              font_size=10.5, header_size=11)

    add_text(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.4),
             "HOW THESE SOURCES MAP TO THE BUSINESS",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.7), Inches(4.6), Inches(12), Inches(2.4), [
        "SuVeechi — answers \"where is each torpedo right now\". Drives the live map and the live trip view.",
        "WBATNGL — answers \"what was tapped, weighed, and dispatched from the BF side\". Drives producer-side reporting.",
        "HTS — answers \"what was received at the converter and what heat was made\". Drives consumer-side reporting.",
        "Each source is owned by a different JSW team (SuVeechi vendor for GPS; BF Operations for WBATNGL; SMS IT for HTS).",
    ], font_size=13)


def slide_d_integration_journey(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "D", "DATA SOURCES & INTEGRATION",
                     "How each source was unblocked — dates and owners",
                     slide_number=n)

    events = [
        ["17-Mar", "HMD Data Requirement Sheet shared with Gagan Chopra"],
        ["23-Mar", "Pritam Saha sets up Torpedo data clarification meeting"],
        ["25-Mar", "Updated Data Requirements re-shared after meeting"],
        ["31-Mar", "Gagan Chopra approves SuVeechi integration path"],
        ["02-Apr", "SuVeechi MySQL view (vw_unit_status_ist) shared by Ganesha (vendor)"],
        ["03-Apr", "Oracle HTS view credentials shared by Hari Prasad (JSW IT)"],
        ["14-Apr", "DEP008 closed — Gagan approval propagated; remaining Oracle accesses cleared"],
        ["16-Apr", "Oracle WBATNGL 3rd view shared by Kotaiah (JSW DBA)"],
        ["06-May", "VPN access granted (DEP012 closed); SuVeechi sync live; WBATNGL sync verified"],
        ["08-May", "WBATNGL trip mirror sprint shipped; per-torpedo capacity backfilled"],
        ["11-May", "HTS connectivity unblocked; correct service name + user identified; DEP007 closed"],
    ]
    rows = [[(d, NAVY, True), e] for d, e in events]
    add_table(s, Inches(0.7), Inches(1.6), Inches(12), Inches(4.6),
              ["Date (2026)", "Event"],
              rows,
              col_widths=[Inches(1.6), Inches(10.4)],
              font_size=11, header_size=12)

    add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.7),
             "Five distinct JSW teams contributed to the integration — Production (Irfan, Gagan), IT Infrastructure (Divakar, Vasagerappa, Pradeep), Oracle DBA (Kotaiah), SMS IT (Hari Prasad), and the SuVeechi vendor (Ganesha, Sanjay). Coordination was the critical-path constraint, not engineering.",
             font_size=11, italic=True, color=GRAY_MED)


def slide_d_provides(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "D", "DATA SOURCES & INTEGRATION",
                     "What each source provides — table view",
                     slide_number=n)

    rows = [
        [("SuVeechi (MySQL)", NAVY, True), "Torpedo identifier (`unitname`)\nLatitude / Longitude\nStatus (Idle / Moving / Ign Off)\nPlant location (free text)\nReport timestamp"],
        [("WBATNGL (Oracle)", NAVY, True), "Trip identifier\nTorpedo (`LADLENO`)\nProducer (`SOURCE_LAB`)\nWeighbridge weights (tare, gross, net)\nTap data (tap_no, tap_hole, temp, S, Si)\nLifecycle timestamps (first_tare, closetime, out_date, sms_ack)"],
        [("HTS (Oracle)", NAVY, True), "Heat identifier (`HEAT_NO`)\nConverter (D/E/F/G/H/I) and SMS unit\nTorpedo arrival/release time\nHot-metal quantity received\nCaster heat process events (SMS-4 only)\nConsumption / yield data (SMS-4 only)"],
    ]
    add_table(s, Inches(0.7), Inches(1.6), Inches(12), Inches(5.0),
              ["Source", "Fields provided"],
              rows,
              col_widths=[Inches(3.0), Inches(9.0)],
              font_size=11, header_size=12)


def slide_d_local_storage(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "D", "DATA SOURCES & INTEGRATION",
                     "Local HMD database — what HMD stores",
                     slide_number=n)

    add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
             "WHY HMD MAINTAINS LOCAL MIRRORS OF EXTERNAL SOURCES",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.7), Inches(1.85), Inches(12), Inches(1.8), [
        "Resilience — if any source database is briefly unreachable, HMD continues to operate from its local mirror.",
        "Performance — HMD-side queries (dashboards, reports) hit local PostgreSQL with sub-second response, not the source.",
        "Reduced load on JSW source systems — HMD pulls deltas at controlled intervals (10s / 60s / 5min) rather than per-user.",
        "Application-owned state stays on HMD side — user accounts, alerts, audit trail, configuration, fleet master, planning data.",
    ], font_size=12)

    add_text(s, Inches(0.7), Inches(3.9), Inches(12), Inches(0.4),
             "WHAT'S STORED LOCALLY",
             font_size=11, bold=True, color=GOLD)
    rows = [
        [("Application state", NAVY, True), "Users, roles, settings, system configuration, alerts, audit log"],
        [("Fleet & nodes", NAVY, True),     "53 torpedoes (master), plant nodes, weighbridges, converters"],
        [("Planning & trips", NAVY, True),  "Daily plans, distribution assignments, trip lifecycle records, route configurations"],
        [("Mirrors of source data", NAVY, True), "fleet_live_locations (SuVeechi snapshot), wbatngl_trip_mirror, hts_heat_mirror + caster mirrors"],
    ]
    add_table(s, Inches(0.7), Inches(4.25), Inches(12), Inches(2.4),
              ["Category", "Contents"],
              rows,
              col_widths=[Inches(2.8), Inches(9.2)],
              font_size=11, header_size=11)


# ─── SECTION E — Application Capabilities ─────────────────────────

def slide_section_e(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(s, "E", "Application Capabilities",
                        "V7 — the production-ready Hot Metal Distribution application.")


def slide_e_modules_overview(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "E", "APPLICATION CAPABILITIES — V7",
                     "Module overview — 12 functional areas",
                     slide_number=n)

    # Grid of 12 modules in 3 rows of 4
    modules = [
        ("Dashboard", "KPI overview, alerts summary"),
        ("Trip Management", "16-state lifecycle, manual & assigned trips"),
        ("Live Tracking", "Real-time torpedo map, status, current trip"),
        ("Operations", "Per-node operator view (Producer / Consumer)"),
        ("Daily Planning", "Capacity, distribution, daily routes"),
        ("Strategic Planning", "Monthly plan, HM Matrix, optimization"),
        ("Torpedo Management", "Fleet registry, status, capacity, lifecycle"),
        ("Maintenance Scheduling", "Calendar view, overlap prevention"),
        ("Statistics", "KPIs, throughput, deviation analytics"),
        ("Reports", "PDF / Excel / CSV exports, 10+ report types"),
        ("Audit Trail", "Username, IP, change diff, regulatory log"),
        ("Settings & Admin", "Users, roles, configurations"),
    ]
    cols, rows_n = 4, 3
    cell_w = Inches(3.05)
    cell_h = Inches(1.55)
    grid_x = Inches(0.55)
    grid_y = Inches(1.55)
    for idx, (title, desc) in enumerate(modules):
        r, c = divmod(idx, cols)
        x = grid_x + c * (cell_w + Inches(0.05))
        y = grid_y + r * (cell_h + Inches(0.15))
        # cell background
        add_solid_rect(s, x, y, cell_w, cell_h, OFF_WHITE)
        # accent stripe
        add_solid_rect(s, x, y, cell_w, Inches(0.06), NAVY)
        # title
        add_text(s, x + Inches(0.15), y + Inches(0.18), cell_w - Inches(0.3), Inches(0.4),
                 title,
                 font_size=14, bold=True, color=NAVY)
        # description
        add_text(s, x + Inches(0.15), y + Inches(0.6), cell_w - Inches(0.3), Inches(0.85),
                 desc,
                 font_size=10, color=GRAY_MED)

    add_text(s, Inches(0.55), Inches(6.55), Inches(12.4), Inches(0.5),
             "All 12 modules are role-aware. Admin sees everything; Producer / Consumer see their node-relevant views only.",
             font_size=11, italic=True, color=GRAY_MED, align=PP_ALIGN.CENTER)


def slide_e_trip_lifecycle(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "E", "APPLICATION CAPABILITIES — V7",
                     "Trip lifecycle management — 16 states from creation to completion",
                     slide_number=n)

    # The lifecycle as horizontal flow with 16 stages
    stages = [
        ("0", "Pending"),
        ("1", "Assigned"),
        ("2", "WB Tare Entry"),
        ("3", "WB Tare Recorded"),
        ("4", "Producer Entered"),
        ("5", "Loading Started"),
        ("6", "Loading Ended"),
        ("7", "Producer Exited"),
        ("8", "WB Gross Entry"),
        ("9", "WB Gross Recorded"),
        ("10", "Consumer Entered"),
        ("11", "Unloading Started"),
        ("12", "Unloading Ended"),
        ("13", "Completed"),
        ("14", "Canceled"),
        ("15", "Aborted"),
    ]
    # Show in 2 rows of 8
    cell_w = Inches(1.5)
    cell_h = Inches(0.9)
    start_x = Inches(0.55)
    start_y = Inches(1.55)
    for idx, (sid, name) in enumerate(stages):
        r, c = divmod(idx, 8)
        x = start_x + c * (cell_w + Inches(0.08))
        y = start_y + r * (cell_h + Inches(0.15))
        # special colour for terminal states
        if sid == "13":
            bg, fg = GREEN, WHITE
        elif sid in ("14", "15"):
            bg, fg = AMBER, WHITE
        else:
            bg, fg = OFF_WHITE, GRAY_DARK
        add_solid_rect(s, x, y, cell_w, cell_h, bg)
        add_text(s, x, y + Inches(0.05), cell_w, Inches(0.3),
                 sid, font_size=10, bold=True, color=fg,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + Inches(0.35), cell_w, Inches(0.5),
                 name, font_size=9, color=fg,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Notes underneath
    add_text(s, Inches(0.55), Inches(3.85), Inches(12.4), Inches(0.4),
             "FEATURES",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(4.2), Inches(12.4), Inches(2.7), [
        "Status transitions are validated server-side — invalid jumps blocked by a state-machine.",
        "Each transition is timestamped and audit-logged with the user who performed it.",
        "Live deviation tracking compares expected times (from HM Matrix) against actuals; threshold alerts fire at 10, 20, 30 minute deltas.",
        "Trips can be created automatically by the assignment engine, manually by Admin, or auto-detected from upstream weighbridge data (WBATNGL mirror).",
        "Trip cancellation is governed by role permissions; abortion mid-execution requires admin override.",
    ], font_size=12)


def slide_e_planning(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "E", "APPLICATION CAPABILITIES — V7",
                     "Planning — daily, monthly, and optimised distribution",
                     slide_number=n)

    # 3 columns
    cols = [
        ("DAILY PLANNING",
         ["Per-day capacity entry by Producer / Consumer roles",
          "Status workflow: Primary → Revised → Confirmed",
          "Auto-computed daily plan rolls into monthly view",
          "Cache-backed for instant operator response"]),
        ("MONTHLY / STRATEGIC PLANNING",
         ["Admin-only monthly capacity calendar",
          "Five planning tabs: Executive, Configuration, Maintenance, Strategic, Weighbridge",
          "Distribution Assignment with Committed / Proposed / Superseded states",
          "Historical comparison and forecast view"]),
        ("HM MATRIX & OPTIMISATION",
         ["Hot-Metal Matrix — bidirectional travel time table per producer-consumer pair",
          "Fill, Unload, Wait time configuration per node",
          "Linear-programming optimiser (PuLP) generates day-optimal allocation",
          "Optimisation considers: capacity, distance, dwell, queue, and torpedo availability"]),
    ]
    col_w = Inches(4.05)
    x0 = Inches(0.55)
    for i, (title, items) in enumerate(cols):
        x = x0 + i * (col_w + Inches(0.1))
        # Top strip
        add_solid_rect(s, x, Inches(1.55), col_w, Inches(0.5), NAVY)
        add_text(s, x + Inches(0.1), Inches(1.55), col_w - Inches(0.2), Inches(0.5),
                 title,
                 font_size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_solid_rect(s, x, Inches(2.05), col_w, Inches(4.5), OFF_WHITE)
        add_bullets(s, x + Inches(0.15), Inches(2.15), col_w - Inches(0.3), Inches(4.3),
                    items, font_size=11)


def slide_e_operations_stats(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "E", "APPLICATION CAPABILITIES — V7",
                     "Operations · Statistics · Deviation Analytics",
                     slide_number=n)

    sections = [
        ("OPERATIONS", "Real-time node-specific control surface",
         ["Per-producer and per-consumer dashboards",
          "Admin sees tabbed view across all producers / consumers",
          "Live status of each node, current torpedo, in-progress trips",
          "Configurable thresholds — warning (10 min), alert (20 min), critical (30 min)",
          "Trip-stage progression with phase-level deviation tracking"]),
        ("STATISTICS",  "Multi-dimensional KPI dashboard",
         ["Producer-specific, Consumer-specific, and global views",
          "Headline KPIs: trips today, tonnes today, average cycle, on-time %",
          "Throughput trends (hourly / daily / monthly)",
          "Fleet utilisation, torpedo cycle distribution, route comparison",
          "All filtered by date range, producer, consumer, or grade"]),
        ("DEVIATION ANALYTICS", "Six analytical lenses on schedule adherence",
         ["Deviation summary by category (Early / On-Time / Warning / Alert / Critical)",
          "By node — per-producer and per-consumer breakdown",
          "By phase — Loading / Transit / Unloading split",
          "Trends — daily and monthly time-series",
          "Period-over-period comparison",
          "Root cause analysis: shift, day-of-week, worst routes"]),
    ]
    y = Inches(1.55)
    for label, sub, items in sections:
        # Header strip
        add_solid_rect(s, Inches(0.55), y, Inches(12.4), Inches(0.5), NAVY_LITE)
        add_text(s, Inches(0.7), y, Inches(3.5), Inches(0.5),
                 label, font_size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.4), y, Inches(8.4), Inches(0.5),
                 sub, font_size=11, italic=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_bullets(s, Inches(0.7), y + Inches(0.55), Inches(12.2), Inches(1.3),
                    items, font_size=11)
        y += Inches(1.85)


def slide_e_tracking_fleet(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "E", "APPLICATION CAPABILITIES — V7",
                     "Live Tracking · Torpedo Management · Maintenance",
                     slide_number=n)

    cols = [
        ("LIVE TRACKING",
         "Real-time map of every torpedo's position",
         ["Leaflet-based interactive map of plant area",
          "53 torpedoes shown with status-coloured markers (Idle / Moving / Ign Off)",
          "Click any torpedo → Torpedo Drawer slides in with:",
          "  • Live latitude / longitude and last update time",
          "  • Current trip details + recent trip history",
          "  • Per-torpedo capacity (from WBATNGL backfill)",
          "  • Maintenance schedule and ladle attributes"]),
        ("TORPEDO MANAGEMENT",
         "Fleet registry — 53 torpedoes",
         ["Master list of every torpedo (TLC-01 through TLC-53)",
          "Status workflow: Operating / Maintenance / Moving",
          "Per-torpedo capacity (range observed: 425–485 MT)",
          "Capacity auto-refreshed nightly from WBATNGL historical weighbridge data",
          "Soft-delete (retired torpedoes preserved for audit)"]),
        ("MAINTENANCE SCHEDULING",
         "Calendar-based maintenance planning",
         ["Schedule maintenance windows per node or torpedo",
          "Overlap prevention to avoid scheduling conflicts",
          "Status tracking: Ongoing / Scheduled / Completed",
          "Modal-based create / edit with role permissions",
          "Duration auto-calculated"]),
    ]
    col_w = Inches(4.05)
    x0 = Inches(0.55)
    for i, (title, sub, items) in enumerate(cols):
        x = x0 + i * (col_w + Inches(0.1))
        add_solid_rect(s, x, Inches(1.55), col_w, Inches(0.85), NAVY)
        add_text(s, x + Inches(0.15), Inches(1.6), col_w - Inches(0.3), Inches(0.4),
                 title, font_size=12, bold=True, color=WHITE)
        add_text(s, x + Inches(0.15), Inches(1.95), col_w - Inches(0.3), Inches(0.4),
                 sub, font_size=10, italic=True,
                 color=RGBColor(0xC7, 0xD2, 0xFE))
        add_solid_rect(s, x, Inches(2.4), col_w, Inches(4.6), OFF_WHITE)
        add_bullets(s, x + Inches(0.15), Inches(2.5), col_w - Inches(0.3), Inches(4.4),
                    items, font_size=10.5)


def slide_e_audit_reports(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "E", "APPLICATION CAPABILITIES — V7",
                     "Reports · Audit Trail · Settings & Security",
                     slide_number=n)

    cols = [
        ("REPORTS",
         ["10+ predefined report types — trip history, planning vs actual, deviation summary, fleet utilisation, throughput, etc.",
          "Export formats: PDF (with branding), Excel, CSV",
          "Filterable by date range, producer, consumer, grade, role",
          "Scheduled reports with history retention",
          "Saved-report configurations for repeat exports"]),
        ("AUDIT TRAIL",
         ["Every action by every user logged",
          "Username, IP address, user agent, action, before/after diff",
          "Filterable, searchable, exportable",
          "Atomic / transaction-safe with retry logic",
          "Regulatory compliance posture"]),
        ("SETTINGS & ADMIN",
         ["User management — create, role-assign, lock, unlock, soft-delete",
          "System configuration — travel times, fill times, wait times, default capacities",
          "Deviation threshold configuration",
          "Notification preferences",
          "Centralised admin control panel"]),
        ("SECURITY POSTURE",
         ["JWT authentication with 480-min token expiry (configurable)",
          "bcrypt password hashing",
          "Account lockout after 5 failed attempts (15-min window)",
          "CSRF protection (double-submit cookie pattern)",
          "Rate limiting (tiered: auth, high, medium, low)",
          "Security headers (X-Content-Type-Options, CSP, etc.)"]),
    ]
    # 2x2 grid
    cell_w = Inches(6.15)
    cell_h = Inches(2.55)
    positions = [
        (Inches(0.55), Inches(1.55)),
        (Inches(6.85), Inches(1.55)),
        (Inches(0.55), Inches(4.25)),
        (Inches(6.85), Inches(4.25)),
    ]
    for (x, y), (title, items) in zip(positions, cols):
        add_solid_rect(s, x, y, cell_w, Inches(0.45), NAVY)
        add_text(s, x + Inches(0.15), y, cell_w - Inches(0.3), Inches(0.45),
                 title, font_size=12, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_solid_rect(s, x, y + Inches(0.45), cell_w, cell_h - Inches(0.45), OFF_WHITE)
        add_bullets(s, x + Inches(0.15), y + Inches(0.55), cell_w - Inches(0.3),
                    cell_h - Inches(0.6), items, font_size=10)


# ─── SECTION F — Current Operations Status ────────────────────────

def slide_section_f(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(s, "F", "Current Operations Status",
                        "What's live today, what data is flowing, and what's been verified end-to-end.")


def slide_f_running(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "F", "CURRENT OPERATIONS STATUS",
                     "What's running on the BF4 PC today",
                     slide_number=n)

    rows = [
        [("HMD backend (FastAPI)", NAVY, True), "Port 8000", ("Running", GREEN, True), "Always-on; service-managed"],
        [("HMD frontend (React)", NAVY, True), "Port 5173 (dev) / built static", ("Running", GREEN, True), "Served to JSW intranet browsers"],
        [("PostgreSQL", NAVY, True), "Port 5432 (localhost)", ("Running", GREEN, True), "Local DB for HMD application state + source mirrors"],
        [("Redis", NAVY, True), "Port 6379 (localhost)", ("Running", GREEN, True), "Cache; auto-fallback to in-memory if unavailable"],
        [("SuVeechi sync job", NAVY, True), "Every 10 seconds", ("Live", GREEN, True), "53 torpedoes — GPS + status snapshot"],
        [("WBATNGL trip sync job", NAVY, True), "Every 60 seconds", ("Live", GREEN, True), "BF-side trip transactions"],
        [("HTS sync job", NAVY, True), "Every 5 minutes", ("Live", GREEN, True), "SMS-side heat records (recently established)"],
        [("WBATNGL capacity backfill", NAVY, True), "Daily at 03:00 IST", ("Live", GREEN, True), "Refreshes per-torpedo capacity from historical net weights"],
    ]
    add_table(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(4.9),
              ["Service / Job", "Frequency", "Status", "Notes"],
              rows,
              col_widths=[Inches(3.0), Inches(2.6), Inches(1.4), Inches(5.4)],
              font_size=11, header_size=12)


def slide_f_data_flowing(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "F", "CURRENT OPERATIONS STATUS",
                     "Live data flowing into HMD",
                     slide_number=n)

    # 3 KPI tiles
    tiles = [
        ("SuVeechi GPS", "53", "torpedoes streaming",
         "Refreshed every 10 seconds.\nDelivers latitude, longitude, ignition status, and free-text plant location for every active torpedo."),
        ("WBATNGL trips", "~3,300", "trip records / 30 days",
         "BF3 + BF5 schemas, union view BF3.WB_TRANS_DATA_ITRO.\nDaily volume: tens to hundreds of trips, depending on production rate."),
        ("HTS heats", "~35,000", "heat records (historical)",
         "Full 10-year history available; current rate steady. SMS-2 and SMS-4 converters represented."),
    ]
    tile_w = Inches(4.05)
    x0 = Inches(0.55)
    for i, (label, value, unit, body) in enumerate(tiles):
        x = x0 + i * (tile_w + Inches(0.1))
        add_solid_rect(s, x, Inches(1.55), tile_w, Inches(2.2), OFF_WHITE)
        add_solid_rect(s, x, Inches(1.55), tile_w, Inches(0.05), GOLD)
        add_text(s, x + Inches(0.2), Inches(1.7), tile_w - Inches(0.4), Inches(0.35),
                 label, font_size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), Inches(2.05), tile_w - Inches(0.4), Inches(0.8),
                 value, font_size=36, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), Inches(2.95), tile_w - Inches(0.4), Inches(0.4),
                 unit, font_size=11, italic=True, color=GRAY_MED)
        add_text(s, x + Inches(0.2), Inches(3.4), tile_w - Inches(0.4), Inches(0.4),
                 body, font_size=10, color=GRAY_DARK)

    # Notes
    add_text(s, Inches(0.55), Inches(4.3), Inches(12.4), Inches(0.4),
             "OBSERVATIONS",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(4.65), Inches(12.4), Inches(2.4), [
        "Data flow is steady. No source has gone offline since the integration went live in early May.",
        "The HMD application's dashboard refreshes automatically as new data arrives — no manual refresh required.",
        "Local mirrors mean operators see sub-second response times even when the underlying source is slow.",
        "Per-torpedo capacities were backfilled from WBATNGL historical data (replacing the 360-MT default with real ranges between 425–485 MT).",
    ], font_size=12)


def slide_f_verified(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "F", "CURRENT OPERATIONS STATUS",
                     "End-to-end verification on the BF4 PC",
                     slide_number=n)

    add_text(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(0.4),
             "VERIFIED FLOWS (AS OF 14-MAY-2026)",
             font_size=11, bold=True, color=GOLD)
    rows = [
        [("Torpedo GPS → live map", NAVY, True), ("Verified", GREEN, True), "Source torpedo position appears on HMD map within ~10 seconds of update"],
        [("Trip dispatched (BF) → HMD record", NAVY, True), ("Verified", GREEN, True), "New row in BF3.WB_TRANS_DATA_ITRO appears in wbatngl_trip_mirror within 60 seconds"],
        [("Heat at SMS → HMD record", NAVY, True), ("Verified", GREEN, True), "New heat in HTS appears in hts_heat_mirror within 5 minutes"],
        [("Per-torpedo capacity refresh", NAVY, True), ("Verified", GREEN, True), "Nightly 03:00 IST job runs; capacities updated from historical net weights"],
        [("Trip lifecycle update → audit trail", NAVY, True), ("Verified", GREEN, True), "Every transition logged with user, IP, before/after"],
        [("Role-based access enforcement", NAVY, True), ("Verified", GREEN, True), "Admin / Producer / Consumer routes correctly restrict views and actions"],
        [("Backup of local DB", NAVY, True), ("Pending", AMBER, True), "Backup schedule to be agreed with JSW IT once go-live is set"],
    ]
    add_table(s, Inches(0.55), Inches(1.9), Inches(12.4), Inches(4.4),
              ["Flow", "Status", "Notes"],
              rows,
              col_widths=[Inches(4.5), Inches(1.6), Inches(6.3)],
              font_size=11, header_size=11)


# ─── SECTION G — Data Reality ─────────────────────────────────────

def slide_section_g(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(s, "G", "Data Reality",
                        "What's reliable, what isn't, and the gaps that must be closed at source.")


def slide_g_module_matrix(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "G", "DATA REALITY",
                     "Module status matrix — what works on current data",
                     slide_number=n)

    rows = [
        [("Trip lifecycle tracking", NAVY, True), ("Functional with gaps", AMBER, True), "Producer side reliable; consumer-side acknowledgement inconsistent"],
        [("Producer-side data (BF / COREX)", NAVY, True), ("Functional with gaps", AMBER, True), "All 7 producers visible; Si chemistry only BF4; S anomalies non-BF4"],
        [("Consumer-side data (all 4 SMS)", NAVY, True), ("Partially blocked", AMBER, True), "SMS-2 and SMS-4 visible; SMS-1 and SMS-3 not represented"],
        [("Fleet management (53 torpedoes)", NAVY, True), ("Functional with inconsistencies", AMBER, True), "Fleet count discrepancy (HMD 53 vs HTS master 36)"],
        [("Live positioning (GPS)", NAVY, True), ("Functional with reliability gaps", AMBER, True), "Coverage complete; vendor-side telemetry limited"],
        [("Geofencing / plant nodes", NAVY, True), ("Partially blocked", RED, True), "No polygons available; node master list informal"],
        [("Weighbridge audit", NAVY, True), ("Functional with gaps", AMBER, True), "3 weighbridges inferred; no calibration history"],
        [("BF-side equipment downtime", NAVY, True), ("Blocked", RED, True), "No source for producer-side downtime events"],
        [("Production planning / targets", NAVY, True), ("Blocked", RED, True), "No source for plans / targets / schedules"],
        [("SMS performance (yield / loss)", NAVY, True), ("Partially blocked", RED, True), "SMS-4 only; SMS-1/2/3 caster data not accessible"],
        [("Heat trace (BF → SMS lineage)", NAVY, True), ("Approximate", AMBER, True), "No shared identifier; time-window join only"],
        [("Alerts & exceptions", NAVY, True), ("Functional", GREEN, True), "Detector running; coverage tracks upstream data availability"],
    ]
    add_table(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(5.5),
              ["Module", "Status", "Reason"],
              rows,
              col_widths=[Inches(4.0), Inches(2.6), Inches(5.8)],
              font_size=10.5, header_size=11)


def slide_g_findings(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "G", "DATA REALITY",
                     "Key data findings — what JSW source data is missing or inconsistent",
                     slide_number=n)

    rows = [
        ["1", "Geofence polygons", "No source provides plant-node geofences. SuVeechi vendor likely has them internally but does not expose. Required for boundary-sensitive events (entry/exit, dwell time)."],
        ["2", "SMS-1 / SMS-3 absent", "HTS contains heat records only for SMS-2 and SMS-4. SMS-1 not represented at all; SMS-3 likely uses EAF (out of scope for hot metal). Status confirmation needed."],
        ["3", "BF3 silent since 24-Sep-2025", "Zero trip activity recorded for BF3 in the past ~232 days. No source explains whether decommissioned, under refractory campaign, or temporarily idle."],
        ["4", "Si chemistry BF4-only", "Of seven producers, only BF4 has silicon (Si) values populated. The other six producers report Si as NULL. Plant-wide Si specification cannot be enforced."],
        ["5", "Sulfur (S) values inconsistent", "BF4 S values fall in plausible 0.015–0.050 band; non-BF4 values reach up to 2.83. Indicates unit or scale inconsistency at source."],
        ["6", "Producer-side downtime missing", "HTS has equipment breakdown data for SMS side. No equivalent table exists or is exposed for any of the 7 producers."],
        ["7", "Production targets unavailable", "No source database contains daily, monthly, or shift-level production targets. Every 'vs plan' comparison is structurally blocked."],
        ["8", "WBATNGL acknowledgement gaps", "SMS-side acknowledgement (SMS_ACK_TIME) is NULL on approximately 45% of trip records. Trip-completion semantics unclear."],
        ["9", "Heat-to-trip linking key absent", "WBATNGL and HTS share only the torpedo identifier. No common HEAT_NO or TRIP_ID exists between the two systems."],
        ["10", "SMS-2 caster data not readable", "Caster process and consumption data for SMS-2 converters (D, E, F, G) appears to be in SPT001A schema. The integration user does not have read access."],
    ]
    add_table(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(5.5),
              ["#", "Finding", "Implication"],
              rows,
              col_widths=[Inches(0.45), Inches(3.5), Inches(8.45)],
              font_size=10, header_size=11)


def slide_g_why_it_matters(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "G", "DATA REALITY",
                     "Operational implications — what these gaps block",
                     slide_number=n)

    rows = [
        ["Production verification", "Producer-side trips visible (all 7 BFs+COREX); consumer-side only for SMS-2/4. Comparing 'produced vs consumed' is feasible for the visible flow but incomplete plant-wide."],
        ["Real-time torpedo location", "GPS streaming is reliable. Boundary detection (torpedo entered SMS-4 at HH:MM) is currently inferential — geofences from SuVeechi vendor would make it deterministic."],
        ["Yield and loss analytics", "Available for SMS-4 only. Approximately 95% of plant heats run on SMS-2 — their yield, loss, and operator notes are not visible. SPT001A access would unlock this."],
        ["Plan-vs-actual reporting", "No targets means no variance reports. The dashboards report absolute throughput but cannot show 'X% of daily plan achieved'."],
        ["Cause-and-effect analysis", "Production drops observed in data have no corresponding cause records. The 'why' must come from operations team interviews rather than source data."],
        ["Plant-wide chemistry compliance", "Si bounds enforceable only on BF4 trips. S anomalies on non-BF4 sources prevent uniform sulfur-spec alerting."],
    ]
    add_table(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(4.9),
              ["Operational area", "Status given current data"],
              rows,
              col_widths=[Inches(3.4), Inches(9.0)],
              font_size=11, header_size=12)


def slide_g_gap_report_ref(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "G", "DATA REALITY",
                     "A formal Data Gap Analysis report has been prepared",
                     slide_number=n)

    add_text(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(0.4),
             "REPORT — HMD ↔ JSW DATA GAP ANALYSIS (V1.0)",
             font_size=11, bold=True, color=GOLD)
    add_text(s, Inches(0.55), Inches(1.95), Inches(12.4), Inches(1.0),
             "A detailed companion document inventories every data requirement of the HMD system, contrasts it with what JSW currently provides across four source systems, and identifies every gap — missing, inconsistent, or fragmented — that prevents the system from operating on reliable, deterministic data.",
             font_size=13)

    rows = [
        [("Coverage", NAVY, True), "13 operational modules across 4 thematic clusters"],
        [("Findings", NAVY, True), "30 distinct data gaps + 5 cross-references"],
        [("Severity breakdown", NAVY, True), "7 BLOCKER · 17 SIGNIFICANT · 5 INCONSISTENT · 1 ENHANCEMENT"],
        [("Per-gap content", NAVY, True), "9-field summary + technical detail + sample data + verification SQL"],
        [("Owner tagging", NAVY, True), "Each gap routed to: JSW IT / JSW Operations / JSW Planning / Vendor"],
        [("Companion document", NAVY, True), "Available as a separate PDF for technical review"],
    ]
    add_table(s, Inches(0.55), Inches(3.0), Inches(12.4), Inches(3.0),
              ["Aspect", "Detail"],
              rows,
              col_widths=[Inches(3.0), Inches(9.4)],
              font_size=11, header_size=11)

    add_text(s, Inches(0.55), Inches(6.3), Inches(12.4), Inches(0.7),
             "The Data Gap Analysis document is grounded in direct schema inspection and concrete query evidence collected during May 2026. It is structured to be acted on per gap by JSW IT, JSW Operations, and JSW Planning.",
             font_size=11, italic=True, color=GRAY_MED)


# ─── SECTION H — Clarifications Needed ────────────────────────────

def slide_section_h(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(s, "H", "Clarifications Needed from JSW",
                        "Specific items where Deevia requires confirmation or input to proceed.")


def slide_h_end_user(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "H", "CLARIFICATIONS NEEDED FROM JSW",
                     "1. End-user identification and engagement",
                     slide_number=n)

    add_text(s, Inches(0.55), Inches(1.6), Inches(12.4), Inches(0.4),
             "WHAT WE NEED",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(1.95), Inches(12.4), Inches(2.2), [
        "Identification of the day-to-day operators of the HMD system at JSW.",
        "A working session with those end users to understand how they perform hot-metal distribution today — current tools, current pain points, current decision flow.",
        "Their input on which HMD features map best to their current operational practice — and which features need adaptation.",
    ], font_size=13)

    add_text(s, Inches(0.55), Inches(4.0), Inches(12.4), Inches(0.4),
             "WHY THIS MATTERS",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(4.35), Inches(12.4), Inches(2.5), [
        "The application is built to specification, but field-level usability depends on operator workflow.",
        "End-user sessions surface details that data alone cannot — for example, what fields operators actually fill in during a shift versus what stays blank, and which decisions are made on-screen versus from a printed sheet.",
        "Early engagement also derisks the pilot phase (M4): a system trained operators can use confidently goes live more smoothly.",
    ], font_size=13)


def slide_h_verification(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "H", "CLARIFICATIONS NEEDED FROM JSW",
                     "2. Production verification approach",
                     slide_number=n)

    add_text(s, Inches(0.55), Inches(1.6), Inches(12.4), Inches(0.4),
             "PROPOSED APPROACH",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(1.95), Inches(12.4), Inches(3.0), [
        "Use trip data from both the producer side (WBATNGL) and the consumer side (HTS) as the primary signal for what was actually produced and consumed during a shift, day, or month.",
        "Producer-side: sum of net-weight tonnes per producer over a period (covering all 7 BFs/COREX).",
        "Consumer-side: sum of hot-metal-quantity tonnes received per converter over the same period (covering SMS-2 and SMS-4 with current access; SMS-1 / SMS-3 / SMS-2 caster pending JSW input).",
        "Any persistent discrepancy between producer and consumer totals would be the operational signal — slag carryback, route losses, weighing drift, or missing nodes.",
    ], font_size=12.5)

    add_text(s, Inches(0.55), Inches(5.2), Inches(12.4), Inches(0.4),
             "CONFIRMATION REQUESTED",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(5.55), Inches(12.4), Inches(1.5), [
        "Is this approach aligned with JSW's current production-reporting methodology, or is there a different reconciliation method in use today?",
        "Who at JSW signs off on the daily / weekly / monthly production numbers, and what data sources do they use today?",
    ], font_size=13)


def slide_h_geofencing(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "H", "CLARIFICATIONS NEEDED FROM JSW",
                     "3. Geofencing — required for the live tracking module to function reliably",
                     slide_number=n)

    add_solid_rect(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(0.7), RED)
    add_text(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(0.7),
             "GEOFENCE POLYGONS ARE A MUST-HAVE — CURRENTLY NOT AVAILABLE FROM ANY SOURCE",
             font_size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.55), Inches(2.5), Inches(12.4), Inches(0.4),
             "WHAT IS NEEDED",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(2.85), Inches(12.4), Inches(2.2), [
        "Polygon coordinates (latitude/longitude vertex list) for every operational plant node — producers (BF1–BF5, COREX1, COREX2), consumers (SMS-1 through SMS-4), all weighbridges, intermediate yards, gates, and pouring stations.",
        "A canonical list of every named plant node in use — including the intermediate locations that appear in SuVeechi tracking data (LRS1, LRS2, SMS2 North PS, HMY1 PCM, BF4 Entry, EY-AVTC Gate, SY-Track Hopper, etc.).",
        "Each node's role classification — producer / consumer / weighbridge / yard / intermediate / gate.",
    ], font_size=12)

    add_text(s, Inches(0.55), Inches(5.2), Inches(12.4), Inches(0.4),
             "WHERE TO GET THESE",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(5.55), Inches(12.4), Inches(1.6), [
        "Most likely already maintained internally by the SuVeechi vendor (the location-text strings in their feed are evidence of internal geofences).",
        "Alternatively, JSW Operations could provide based on plant layout drawings — even rough polygons would be a significant improvement over single-point coordinates.",
        "Once received, geofences will be loaded into HMD and used to derive deterministic entry/exit/dwell events for every torpedo.",
    ], font_size=12)


# ─── SECTION I — Path Forward ─────────────────────────────────────

def slide_section_i(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_divider(s, "I", "Path Forward",
                        "Open dependencies, next steps, and the completion forecast.")


def slide_i_pending(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "I", "PATH FORWARD",
                     "Pending dependencies and next steps",
                     slide_number=n)

    rows = [
        [("DEP002", NAVY, True), "Level 1/2 production signals + Production Plan access",
         "JSW Operations", ("Open · Overdue", RED, True), "Targeted: 30 days from this review"],
        [("DEP005", NAVY, True), "Production / Consumption rules documentation",
         "JSW Operations", ("Open · Overdue", RED, True), "Targeted: 14 days from this review"],
        [("DEP010", NAVY, True), "BF5.WB_TRANS_DATA_ITRO Oracle view recompile",
         "JSW IT (Kotaiah)", ("In Progress", AMBER, True), "Workaround in place; permanent fix pending"],
        [("New",   NAVY, True), "Geofence polygons for all plant nodes",
         "SuVeechi vendor (via JSW)", ("New ask", GOLD, True), "Targeted: 30 days from this review"],
        [("New",   NAVY, True), "SMS-2 caster data access (SPT001A schema)",
         "JSW IT", ("New ask", GOLD, True), "Targeted: 30 days from this review"],
        [("New",   NAVY, True), "Production targets / heat schedules exposure",
         "JSW Planning + JSW IT", ("New ask", GOLD, True), "Targeted: 30 days from this review"],
        [("New",   NAVY, True), "BF-side equipment downtime source",
         "JSW IT + JSW Operations", ("New ask", GOLD, True), "Targeted: 30 days from this review"],
        [("New",   NAVY, True), "End-user identification + working session",
         "JSW Operations", ("New ask", GOLD, True), "Within 14 days of this review — prerequisite for pilot phase"],
    ]
    add_table(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(5.0),
              ["ID", "Item", "Owner", "Status", "Timeline"],
              rows,
              col_widths=[Inches(1.0), Inches(4.6), Inches(2.8), Inches(1.7), Inches(2.3)],
              font_size=10, header_size=11)


def slide_i_forecast(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "I", "PATH FORWARD",
                     "Completion forecast — M3, M4, M5 scenarios",
                     slide_number=n)

    rows = [
        [("M3", NAVY, True), "Dashboard & Integration",
         "19-Jun-2026", "Same — if pending dependencies (DEP002, DEP005, geofence access) close within 30 days",
         "Mid-Jul-2026 — if dependencies slip"],
        [("M4", NAVY, True), "Testing & Pilot",
         "24-Jul-2026", "Same — assuming M3 closes on plan and end-user availability confirmed",
         "Aug-2026 — under slip scenario"],
        [("M5", NAVY, True), "Go-Live & Stabilization",
         "07-Aug-2026", "Same — pending successful pilot in M4",
         "Aug-Sep 2026 — under slip scenario"],
    ]
    add_table(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(2.8),
              ["ID", "Milestone", "Plan End", "Best-case", "Likely if delayed"],
              rows,
              col_widths=[Inches(0.7), Inches(2.6), Inches(1.5), Inches(4.0), Inches(3.6)],
              font_size=10.5, header_size=11)

    add_text(s, Inches(0.55), Inches(4.55), Inches(12.4), Inches(0.4),
             "KEY MESSAGES",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(4.9), Inches(12.4), Inches(2.1), [
        "Engineering work for V7 is essentially complete. Remaining timeline is JSW-dependent — data exposure, operational confirmations, and end-user engagement.",
        "If the eight pending items in the prior slide close on the targets shown, the original 07-Aug-2026 go-live remains achievable.",
        "Each week of delay on JSW-side items shifts the go-live by approximately one week.",
        "Deevia continues to maintain weekly status reports and is available for accelerated cadence (daily standups) during the pilot phase.",
    ], font_size=12)


# ─── SECTION J — Close ────────────────────────────────────────────

def slide_j_close(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(s, "J", "CLOSE",
                     "Summary and discussion",
                     slide_number=n)

    add_text(s, Inches(0.55), Inches(1.55), Inches(12.4), Inches(0.4),
             "PROJECT SUMMARY",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(1.9), Inches(12.4), Inches(2.0), [
        "HMD V7 — full-featured hot metal distribution application — is built and operational on JSW infrastructure.",
        "All three external data sources are integrated; data flows live every 10 seconds to 5 minutes.",
        "Project tracking at approximately 53% completion overall, on plan for the 07-Aug-2026 go-live.",
        "Outstanding items are predominantly JSW-side: data exposure, operational confirmations, and end-user engagement.",
    ], font_size=13)

    add_text(s, Inches(0.55), Inches(4.1), Inches(12.4), Inches(0.4),
             "WHAT WE NEED FROM JSW",
             font_size=11, bold=True, color=GOLD)
    add_bullets(s, Inches(0.55), Inches(4.45), Inches(12.4), Inches(2.0), [
        "End-user identification and a working session — within 14 days.",
        "Geofence polygons for all plant nodes — within 30 days.",
        "Production targets / heat schedules / production-consumption rules — within 30 days.",
        "SMS-2 caster data access (SPT001A) — within 30 days.",
        "BF-side equipment downtime source — within 30 days.",
    ], font_size=13)

    add_text(s, Inches(0.55), Inches(6.7), Inches(12.4), Inches(0.4),
             "DISCUSSION & Q&A",
             font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ─── Main builder ─────────────────────────────────────────────────

def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1 — cover
    slide_cover(prs)

    # Section A — Project Foundation (slides 2-5)
    slide_agenda(prs, 2)
    slide_project_goal(prs, 3)
    slide_architecture(prs, 4)

    # Section B — Timeline (slides 5-8)
    slide_section_b(prs)
    slide_b_timeline_plan(prs, 6)
    slide_b_timeline_actual(prs, 7)
    slide_b_milestone_progress(prs, 8)

    # Section C — Infrastructure (slides 9-12)
    slide_section_c(prs)
    slide_c_server(prs, 10)
    slide_c_network(prs, 11)
    slide_c_stack(prs, 12)

    # Section D — Data Sources (slides 13-17)
    slide_section_d(prs)
    slide_d_sources_overview(prs, 14)
    slide_d_integration_journey(prs, 15)
    slide_d_provides(prs, 16)
    slide_d_local_storage(prs, 17)

    # Section E — Application Capabilities (slides 18-24)
    slide_section_e(prs)
    slide_e_modules_overview(prs, 19)
    slide_e_trip_lifecycle(prs, 20)
    slide_e_planning(prs, 21)
    slide_e_operations_stats(prs, 22)
    slide_e_tracking_fleet(prs, 23)
    slide_e_audit_reports(prs, 24)

    # Section F — Current Operations (slides 25-28)
    slide_section_f(prs)
    slide_f_running(prs, 26)
    slide_f_data_flowing(prs, 27)
    slide_f_verified(prs, 28)

    # Section G — Data Reality (slides 29-33)
    slide_section_g(prs)
    slide_g_module_matrix(prs, 30)
    slide_g_findings(prs, 31)
    slide_g_why_it_matters(prs, 32)
    slide_g_gap_report_ref(prs, 33)

    # Section H — Clarifications (slides 34-37)
    slide_section_h(prs)
    slide_h_end_user(prs, 35)
    slide_h_verification(prs, 36)
    slide_h_geofencing(prs, 37)

    # Section I — Path Forward (slides 38-40)
    slide_section_i(prs)
    slide_i_pending(prs, 39)
    slide_i_forecast(prs, 40)

    # Section J — Close (slide 41)
    slide_j_close(prs, 41)

    out = Path(__file__).parent / "2026-05-14-hmd-project-status.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    p = build_deck()
    print(f"Generated: {p}")
    print(f"Size: {p.stat().st_size / 1024:.1f} KB")
